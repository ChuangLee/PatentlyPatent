"""v0.37 P0-3: 二进制文档文本提取（PDF / pptx / docx）。

策略：
  - 按 mime 选提取器：pdfplumber / python-pptx / python-docx
  - 失败兜底返回 None（调用方决定怎么呈现）
  - 大文件截断到 200_000 字（一份 300 页 PDF 大约 50k-100k 字，够 LLM 消化）
"""
from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

MAX_CHARS = 200_000

PDF_MIMES = {"application/pdf"}
PPTX_MIMES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}
DOCX_MIMES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
}
XLSX_MIMES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
}
TEXT_MIMES = {"text/plain", "text/csv", "text/markdown", "application/json"}


def can_extract(mime: str | None, name: str | None = None) -> bool:
    """是否支持文本提取（不调用就先判断）。"""
    m = (mime or "").lower()
    if m in PDF_MIMES or m in PPTX_MIMES or m in DOCX_MIMES or m in XLSX_MIMES or m in TEXT_MIMES:
        return True
    # 按文件后缀兜底（部分上传 mime 为 octet-stream）
    lower = (name or "").lower()
    return lower.endswith((".pdf", ".pptx", ".ppt", ".docx", ".xlsx", ".xls",
                            ".txt", ".csv", ".md", ".json"))


def extract_text(path: Path, mime: str | None, name: str | None = None) -> str | None:
    """读 path 文件、提取文本；失败返回 None。"""
    if not path.exists():
        return None
    m = (mime or "").lower()
    lower = (name or path.name).lower()
    try:
        if m in PDF_MIMES or lower.endswith(".pdf"):
            return _extract_pdf(path)
        if m in PPTX_MIMES or lower.endswith((".pptx",)):
            return _extract_pptx(path)
        if lower.endswith(".ppt"):
            return f"_暂不支持 .ppt 老格式（{path.name}），请用户另存为 .pptx 后重新上传。_"
        if m in DOCX_MIMES or lower.endswith(".docx"):
            return _extract_docx(path)
        if lower.endswith(".doc"):
            return f"_暂不支持 .doc 老格式（{path.name}），请用户另存为 .docx 后重新上传。_"
        if lower.endswith(".xlsx"):
            return _extract_xlsx(path)
        if lower.endswith(".xls"):
            return _extract_xls(path)
        if m in TEXT_MIMES or lower.endswith((".txt", ".csv", ".md", ".json")):
            return _truncate(path.read_text(encoding="utf-8", errors="ignore"))
    except Exception as exc:  # noqa: BLE001
        logger.warning("extract_text(%s) failed: %s", path, exc)
        return None
    return None


def _truncate(text: str) -> str:
    if not text:
        return ""
    if len(text) <= MAX_CHARS:
        return text
    return text[:MAX_CHARS] + f"\n\n…[本文已截断到前 {MAX_CHARS} 字（原文长度 {len(text)} 字）]"


def _extract_pdf(path: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            try:
                txt = page.extract_text() or ""
            except Exception:  # noqa: BLE001
                txt = ""
            if txt.strip():
                parts.append(f"\n--- Page {i} ---\n{txt}")
    return _truncate("\n".join(parts).strip())


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_lines.append(text)
        # 备注
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                slide_lines.append(f"[备注] {notes}")
        if slide_lines:
            parts.append(f"\n--- Slide {i} ---\n" + "\n".join(slide_lines))
    return _truncate("\n".join(parts).strip())


def _extract_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    # 正文段落
    for para in doc.paragraphs:
        t = (para.text or "").strip()
        if t:
            parts.append(t)
    # 表格
    for tbl in doc.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return _truncate("\n".join(parts).strip())


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"\n--- Sheet: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append(" | ".join(cells))
    wb.close()
    return _truncate("\n".join(parts).strip())


def _extract_xls(path: Path) -> str:
    import xlrd

    wb = xlrd.open_workbook(str(path))
    parts: list[str] = []
    for sheet in wb.sheets():
        parts.append(f"\n--- Sheet: {sheet.name} ---")
        for r in range(sheet.nrows):
            cells = [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
            if any(cells):
                parts.append(" | ".join(cells))
    return _truncate("\n".join(parts).strip())
